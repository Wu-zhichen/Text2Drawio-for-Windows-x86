from __future__ import annotations

import base64
import csv
import io
import json
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, List, Mapping, Sequence, Tuple


SUPPORTED_ATTACHMENT_EXTENSIONS = (
    ".pdf",
    ".xlsx",
    ".xlsm",
    ".xltx",
    ".xltm",
    ".xls",
    ".csv",
    ".tsv",
    ".docx",
    ".pptx",
    ".txt",
    ".md",
    ".json",
    ".xml",
)

_OFFICE_ZIP_EXTENSIONS = {".xlsx", ".xlsm", ".xltx", ".xltm", ".docx", ".pptx"}


class AttachmentError(ValueError):
    pass


@dataclass(frozen=True)
class ExtractedDocuments:
    context: str
    warnings: List[str]
    filenames: List[str]


class DocumentExtractor:
    """Extract bounded, text-only context from local user-selected files in memory."""

    def __init__(
        self,
        max_files: int = 5,
        max_file_bytes: int = 15 * 1024 * 1024,
        max_chars_per_file: int = 60_000,
        max_total_chars: int = 120_000,
        max_office_uncompressed_bytes: int = 120 * 1024 * 1024,
    ) -> None:
        self.max_files = max(1, max_files)
        self.max_file_bytes = max(1_000_000, max_file_bytes)
        self.max_chars_per_file = max(2_000, max_chars_per_file)
        self.max_total_chars = max(self.max_chars_per_file, max_total_chars)
        self.max_office_uncompressed_bytes = max(
            self.max_file_bytes, max_office_uncompressed_bytes
        )

    def extract(self, attachments: Sequence[Mapping[str, Any]]) -> ExtractedDocuments:
        if not attachments:
            return ExtractedDocuments(context="", warnings=[], filenames=[])
        if len(attachments) > self.max_files:
            raise AttachmentError(f"A maximum of {self.max_files} files can be attached.")

        sections: List[str] = []
        warnings: List[str] = []
        filenames: List[str] = []
        remaining = self.max_total_chars
        for attachment in attachments:
            filename = self._safe_filename(attachment.get("filename"))
            extension = Path(filename).suffix.lower()
            if extension not in SUPPORTED_ATTACHMENT_EXTENSIONS:
                supported = ", ".join(SUPPORTED_ATTACHMENT_EXTENSIONS)
                raise AttachmentError(
                    f"Unsupported attachment type for {filename!r}. Supported: {supported}"
                )
            raw = self._decode(attachment.get("data_base64"), filename)
            if len(raw) > self.max_file_bytes:
                limit_mb = self.max_file_bytes // (1024 * 1024)
                raise AttachmentError(f"{filename!r} exceeds the {limit_mb} MB file limit.")
            if extension in _OFFICE_ZIP_EXTENSIONS:
                self._validate_office_archive(raw, filename)

            text, extraction_notes = self._extract_one(raw, extension, filename)
            text = self._clean_text(text)
            if not text:
                if extension == ".pdf":
                    raise AttachmentError(
                        f"{filename!r} contains no extractable text. It may be a scanned PDF; OCR is not enabled."
                    )
                raise AttachmentError(f"{filename!r} contains no readable text or cells.")

            allowed = min(self.max_chars_per_file, remaining)
            if allowed <= 0:
                warnings.append(
                    f"Skipped {filename} because the total attachment context limit was reached."
                )
                continue
            if len(text) > allowed:
                text = text[:allowed].rstrip()
                warnings.append(f"Truncated extracted text from {filename} to {allowed} characters.")
            warnings.extend(extraction_notes)
            sections.append(
                f"<<<ATTACHMENT name={json.dumps(filename, ensure_ascii=False)} type={extension[1:]}>>>\n"
                f"{text}\n<<<END ATTACHMENT>>>"
            )
            filenames.append(filename)
            remaining -= len(text)

        if not sections:
            raise AttachmentError("No usable attachment content was extracted.")
        return ExtractedDocuments(
            context="\n\n".join(sections), warnings=warnings, filenames=filenames
        )

    @staticmethod
    def _safe_filename(value: Any) -> str:
        filename = Path(str(value or "").replace("\\", "/")).name.strip()
        if not filename or filename in {".", ".."} or "\x00" in filename:
            raise AttachmentError("Attachment filename is missing or invalid.")
        return filename[:255]

    @staticmethod
    def _decode(value: Any, filename: str) -> bytes:
        if not isinstance(value, str) or not value:
            raise AttachmentError(f"{filename!r} has no file data.")
        try:
            return base64.b64decode(value, validate=True)
        except (ValueError, TypeError) as exc:
            raise AttachmentError(f"{filename!r} contains invalid Base64 file data.") from exc

    def _validate_office_archive(self, raw: bytes, filename: str) -> None:
        try:
            with zipfile.ZipFile(io.BytesIO(raw)) as archive:
                infos = archive.infolist()
                if len(infos) > 5_000:
                    raise AttachmentError(f"{filename!r} contains too many archive entries.")
                unpacked = sum(item.file_size for item in infos)
                if unpacked > self.max_office_uncompressed_bytes:
                    raise AttachmentError(
                        f"{filename!r} expands beyond the safe Office document limit."
                    )
        except zipfile.BadZipFile as exc:
            raise AttachmentError(f"{filename!r} is not a valid Office document.") from exc

    def _extract_one(
        self, raw: bytes, extension: str, filename: str
    ) -> Tuple[str, List[str]]:
        try:
            if extension == ".pdf":
                return self._pdf(raw), []
            if extension in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
                return self._xlsx(raw), []
            if extension == ".xls":
                return self._xls(raw), []
            if extension in {".csv", ".tsv"}:
                return self._delimited(raw, extension), []
            if extension == ".docx":
                return self._docx(raw), []
            if extension == ".pptx":
                return self._pptx(raw), []
            return self._decode_text(raw), []
        except AttachmentError:
            raise
        except Exception as exc:  # Libraries expose many format-specific exception types.
            raise AttachmentError(f"Could not read {filename!r}: {type(exc).__name__}") from exc

    @staticmethod
    def _pdf(raw: bytes) -> str:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(raw))
        if reader.is_encrypted and not reader.decrypt(""):
            raise AttachmentError("The selected PDF is password-protected.")
        pages: List[str] = []
        for index, page in enumerate(reader.pages[:300], start=1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"[Page {index}]\n{text.strip()}")
        return "\n\n".join(pages)

    @staticmethod
    def _xlsx(raw: bytes) -> str:
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        sections: List[str] = []
        try:
            for worksheet in workbook.worksheets[:30]:
                rows: List[str] = []
                for row_index, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                    values = [DocumentExtractor._cell_text(value) for value in row]
                    while values and not values[-1]:
                        values.pop()
                    if any(values):
                        rows.append(f"{row_index}\t" + "\t".join(values))
                    if sum(len(item) for item in rows) > 80_000:
                        break
                if rows:
                    sections.append(f"[Sheet: {worksheet.title}]\n" + "\n".join(rows))
        finally:
            workbook.close()
        return "\n\n".join(sections)

    @staticmethod
    def _xls(raw: bytes) -> str:
        import xlrd

        workbook = xlrd.open_workbook(file_contents=raw, on_demand=True)
        sections: List[str] = []
        try:
            for worksheet in workbook.sheets()[:30]:
                rows: List[str] = []
                for row_index in range(min(worksheet.nrows, 10_000)):
                    values = [
                        DocumentExtractor._cell_text(worksheet.cell_value(row_index, col_index))
                        for col_index in range(worksheet.ncols)
                    ]
                    while values and not values[-1]:
                        values.pop()
                    if any(values):
                        rows.append(f"{row_index + 1}\t" + "\t".join(values))
                    if sum(len(item) for item in rows) > 80_000:
                        break
                if rows:
                    sections.append(f"[Sheet: {worksheet.name}]\n" + "\n".join(rows))
        finally:
            workbook.release_resources()
        return "\n\n".join(sections)

    @staticmethod
    def _delimited(raw: bytes, extension: str) -> str:
        text = DocumentExtractor._decode_text(raw)
        delimiter = "\t" if extension == ".tsv" else ","
        if extension == ".csv":
            try:
                delimiter = csv.Sniffer().sniff(text[:4096], delimiters=",;\t|").delimiter
            except csv.Error:
                pass
        reader = csv.reader(io.StringIO(text), delimiter=delimiter)
        rows = []
        for index, row in enumerate(reader, start=1):
            if index > 10_000:
                break
            values = [item.strip() for item in row]
            if any(values):
                rows.append(f"{index}\t" + "\t".join(values))
        return "\n".join(rows)

    @staticmethod
    def _docx(raw: bytes) -> str:
        from docx import Document

        document = Document(io.BytesIO(raw))
        parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        for table_index, table in enumerate(document.tables, start=1):
            rows = []
            for row in table.rows:
                values = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                if any(values):
                    rows.append("\t".join(values))
            if rows:
                parts.append(f"[Table {table_index}]\n" + "\n".join(rows))
        return "\n\n".join(parts)

    @staticmethod
    def _pptx(raw: bytes) -> str:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(raw))
        slides = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    texts.append(text.strip())
            if texts:
                slides.append(f"[Slide {index}]\n" + "\n".join(texts))
        return "\n\n".join(slides)

    @staticmethod
    def _decode_text(raw: bytes) -> str:
        for encoding in ("utf-8-sig", "gb18030", "utf-16", "latin-1"):
            try:
                return raw.decode(encoding)
            except UnicodeDecodeError:
                continue
        return raw.decode("utf-8", errors="replace")

    @staticmethod
    def _cell_text(value: Any) -> str:
        if value is None:
            return ""
        if isinstance(value, float) and value.is_integer():
            return str(int(value))
        return str(value).strip().replace("\n", " ")

    @staticmethod
    def _clean_text(text: str) -> str:
        cleaned = text.replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n")
        lines = [line.rstrip() for line in cleaned.split("\n")]
        compact: List[str] = []
        blank = False
        for line in lines:
            if line.strip():
                compact.append(line)
                blank = False
            elif not blank:
                compact.append("")
                blank = True
        return "\n".join(compact).strip()
